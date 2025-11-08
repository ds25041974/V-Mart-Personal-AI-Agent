"""
PowerPoint Reader
Reads data from PowerPoint presentations including:
- All slides (visible and hidden)
- Slide content (text, tables, charts)
- Speaker notes
- Slide layouts and masters
- Embedded objects
- Metadata
"""

import os
import platform
import subprocess
from typing import Any, Dict, List, Optional

from pptx import Presentation


class PowerPointReader:
    """Read data from PowerPoint presentations"""

    def __init__(self, file_path: Optional[str] = None):
        self.file_path = file_path
        self.presentation = None

    def get_active_powerpoint_file(self) -> Optional[str]:
        """Get the path of the currently active PowerPoint file"""
        system = platform.system()

        if system == "Darwin":  # macOS
            return self._get_active_ppt_macos()
        elif system == "Windows":
            return self._get_active_ppt_windows()

        return None

    def _get_active_ppt_macos(self) -> Optional[str]:
        """Get active PowerPoint file path on macOS"""
        try:
            script = """
            tell application "Microsoft PowerPoint"
                if (count of presentations) > 0 then
                    set activePresentation to active presentation
                    set presPath to full name of activePresentation
                    return presPath
                else
                    return ""
                end if
            end tell
            """

            result = subprocess.run(
                ["osascript", "-e", script], capture_output=True, text=True, timeout=5
            )

            if result.returncode == 0 and result.stdout.strip():
                return result.stdout.strip()
        except Exception as e:
            print(f"Error getting active PowerPoint file: {e}")

        return None

    def _get_active_ppt_windows(self) -> Optional[str]:
        """Get active PowerPoint file path on Windows"""
        try:
            script = """
            $ppt = [System.Runtime.InteropServices.Marshal]::GetActiveObject("PowerPoint.Application")
            if ($ppt.ActivePresentation) {
                $ppt.ActivePresentation.FullName
            }
            """

            result = subprocess.run(
                ["powershell", "-Command", script],
                capture_output=True,
                text=True,
                timeout=5,
            )

            if result.returncode == 0 and result.stdout.strip():
                return result.stdout.strip()
        except Exception as e:
            print(f"Error getting active PowerPoint file: {e}")

        return None

    def load_file(self, file_path: Optional[str] = None) -> bool:
        """Load PowerPoint file for reading"""
        if file_path:
            self.file_path = file_path

        if not self.file_path:
            self.file_path = self.get_active_powerpoint_file()

        if not self.file_path or not os.path.exists(self.file_path):
            return False

        try:
            self.presentation = Presentation(self.file_path)
            return True
        except Exception as e:
            print(f"Error loading PowerPoint file: {e}")
            return False

    def get_file_metadata(self) -> Dict[str, Any]:
        """Get PowerPoint file metadata"""
        if not self.presentation:
            return {}

        props = self.presentation.core_properties

        return {
            "filename": os.path.basename(self.file_path) if self.file_path else "",
            "full_path": self.file_path or "",
            "title": props.title or "",
            "author": props.author or "",
            "subject": props.subject or "",
            "keywords": props.keywords or "",
            "created": props.created.isoformat() if props.created else "",
            "modified": props.modified.isoformat() if props.modified else "",
            "last_modified_by": props.last_modified_by or "",
            "revision": props.revision,
            "slide_count": len(self.presentation.slides),
            "slide_width": self.presentation.slide_width,
            "slide_height": self.presentation.slide_height,
        }

    def get_all_slides_info(self) -> List[Dict[str, Any]]:
        """Get information about all slides"""
        if not self.presentation:
            return []

        slides_info = []

        for idx, slide in enumerate(self.presentation.slides, start=1):
            slides_info.append(
                {
                    "slide_number": idx,
                    "slide_id": slide.slide_id,
                    "layout_name": slide.slide_layout.name
                    if slide.slide_layout
                    else "Unknown",
                    "has_notes": bool(
                        slide.notes_slide.notes_text_frame.text.strip()
                        if slide.has_notes_slide
                        else False
                    ),
                    "shape_count": len(slide.shapes),
                }
            )

        return slides_info

    def read_slide_content(self, slide_number: int) -> Dict[str, Any]:
        """Read comprehensive data from a specific slide"""
        if (
            not self.presentation
            or slide_number < 1
            or slide_number > len(self.presentation.slides)
        ):
            return {"error": "Invalid slide number"}

        slide = self.presentation.slides[slide_number - 1]

        # Extract text from shapes
        text_content = []
        tables = []
        charts = []
        images = []

        for shape in slide.shapes:
            # Text boxes and placeholders
            if hasattr(shape, "text") and shape.text:
                text_content.append(
                    {
                        "type": shape.shape_type.name
                        if hasattr(shape.shape_type, "name")
                        else str(shape.shape_type),
                        "text": shape.text,
                        "name": shape.name,
                    }
                )

            # Tables
            if shape.shape_type == 19:  # TABLE
                table_data = self._extract_table_data(shape.table)
                tables.append(
                    {
                        "name": shape.name,
                        "rows": len(shape.table.rows),
                        "columns": len(shape.table.columns),
                        "data": table_data,
                    }
                )

            # Charts
            if hasattr(shape, "has_chart") and shape.has_chart:
                charts.append(
                    {
                        "name": shape.name,
                        "chart_type": shape.chart.chart_type.name
                        if hasattr(shape.chart.chart_type, "name")
                        else str(shape.chart.chart_type),
                        "title": shape.chart.chart_title.text_frame.text
                        if shape.chart.has_title
                        else "",
                    }
                )

            # Images
            if hasattr(shape, "image"):
                images.append(
                    {"name": shape.name, "width": shape.width, "height": shape.height}
                )

        # Get speaker notes
        notes = ""
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            notes = (
                notes_slide.notes_text_frame.text
                if notes_slide.notes_text_frame
                else ""
            )

        return {
            "slide_number": slide_number,
            "slide_id": slide.slide_id,
            "layout_name": slide.slide_layout.name if slide.slide_layout else "Unknown",
            "text_content": text_content,
            "tables": tables,
            "charts": charts,
            "images": images,
            "notes": notes,
            "shape_count": len(slide.shapes),
        }

    def _extract_table_data(self, table) -> List[List[str]]:
        """Extract data from a PowerPoint table"""
        table_data = []

        for row in table.rows:
            row_data = []
            for cell in row.cells:
                row_data.append(cell.text)
            table_data.append(row_data)

        return table_data

    def read_all_slides(self) -> Dict[str, Any]:
        """Read all slides from the presentation"""
        if not self.presentation:
            if not self.load_file():
                return {"error": "Failed to load PowerPoint file"}

        result = {
            "metadata": self.get_file_metadata(),
            "slides_summary": self.get_all_slides_info(),
            "slides_data": [],
        }

        for slide_num in range(1, len(self.presentation.slides) + 1):
            result["slides_data"].append(self.read_slide_content(slide_num))

        return result
